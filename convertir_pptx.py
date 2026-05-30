# -*- coding: utf-8 -*-
"""Renderiza cada slide del HTML a PNG 1920x1080 y arma un PPTX (imagen por slide)."""
import os
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu

BASE = os.path.dirname(os.path.abspath(__file__))
HTML = os.path.join(BASE, "Diapositivas_v2.html")
OUTDIR = os.path.join(BASE, "_slides_png")
PPTX = os.path.join(BASE, "Diapositivas_v2.pptx")
W, H = 1920, 1080

os.makedirs(OUTDIR, exist_ok=True)

def render():
    pngs = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch(channel="chrome", headless=True)
        page = browser.new_page(viewport={"width": W, "height": H}, device_scale_factor=1)
        page.goto("file:///" + HTML.replace("\\", "/"))
        # esperar fuentes/iconos
        page.wait_for_timeout(2500)
        n = page.evaluate("document.querySelectorAll('.slide').length")
        print("slides detectadas:", n)
        # activar modo presentacion con escala fija = 1 (lienzo exacto 1920x1080)
        page.evaluate("""() => {
            document.body.classList.add('pres');
            document.body.style.setProperty('--pres-scale', 1);
            document.querySelectorAll('.slide').forEach(s => s.classList.remove('active','from-left','from-right'));
            // ocultar controles de navegacion y barra de progreso
            ['nav','pbar'].forEach(id => { const e = document.getElementById(id); if (e) e.style.display = 'none'; });
        }""")
        for i in range(n):
            page.evaluate("""(idx) => {
                const s = document.querySelectorAll('.slide');
                s.forEach(x => x.classList.remove('active','from-left','from-right'));
                s[idx].classList.add('active');
            }""", i)
            page.wait_for_timeout(900)  # dejar asentar animaciones
            out = os.path.join(OUTDIR, f"slide_{i:02d}.png")
            page.screenshot(path=out, clip={"x": 0, "y": 0, "width": W, "height": H})
            pngs.append(out)
            print("  ->", os.path.basename(out))
        browser.close()
    return pngs

def build_pptx(pngs):
    prs = Presentation()
    # 16:9 a partir de 1920x1080 (EMU: 914400 por pulgada, 96 dpi -> 1920px = 20in)
    prs.slide_width = Emu(int(W / 96 * 914400))
    prs.slide_height = Emu(int(H / 96 * 914400))
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX)
    print("PPTX guardado en:", PPTX)

if __name__ == "__main__":
    pngs = render()
    build_pptx(pngs)
